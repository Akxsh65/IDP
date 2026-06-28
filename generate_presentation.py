"""
Generate presentation slides for ML framework.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

PROJECT_DIR = Path(__file__).resolve().parent
FIGURES_DIR = PROJECT_DIR / "outputs_publication_figures"
OUTPUT_FILE = PROJECT_DIR / "ML_Framework_Presentation.pptx"

def add_title_slide(prs, title, subtitle=""):
    """Add title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    return slide

def add_content_slide(prs, title, bullets, image_path=None, image_width=5.5):
    """Add content slide with title, bullets, and optional image."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add bullets to left side
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(14)
    
    # Add image to right side if provided
    if image_path and image_path.exists():
        left = Inches(5.2)
        top = Inches(1.5)
        pic = slide.shapes.add_picture(str(image_path), left, top, width=Inches(image_width))
    
    return slide

def add_two_image_slide(prs, title, bullets, image1_path, image2_path, img_width=2.3):
    """Add slide with title, bullets, and two images side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add bullets
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(12)
    
    # Add two images
    if image1_path.exists():
        pic1 = slide.shapes.add_picture(str(image1_path), Inches(0.5), Inches(3.5), width=Inches(img_width))
    
    if image2_path.exists():
        pic2 = slide.shapes.add_picture(str(image2_path), Inches(3.2), Inches(3.5), width=Inches(img_width))
    
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "Machine Learning Framework",
        "Surrogate Modeling for Piezoelectric Morphing Wings"
    )
    
    # Slide 2: ML Problem and Goal
    add_content_slide(
        prs,
        "ML Problem and Goal",
        [
            "Input: Lower MFC voltage, Upper MFC voltage, Pitch angle",
            "Output: Trailing-edge deflection, Lift coefficient (C_L), Drag coefficient (C_D)",
            "Challenge: Replace expensive CFD and structural simulations",
            "Solution: Build fast, accurate surrogate models using ML"
        ],
        image_path=FIGURES_DIR / "fig01_pipeline_schematic.png",
        image_width=4.5
    )
    
    # Slide 3: Data and Training Space
    add_content_slide(
        prs,
        "Data and Training Space",
        [
            "Structural dataset: 300 ANSYS coupled-field simulations",
            "Aerodynamic dataset: 201 CFD samples",
            "Sparse coverage: Limited voltage combinations and pitch angles",
            "Solution: Physics-consistent synthetic data augmentation"
        ],
        image_path=FIGURES_DIR / "fig10_data_coverage.png",
        image_width=4.5
    )
    
    # Slide 4: Structural Surrogate Model
    add_two_image_slide(
        prs,
        "Structural Surrogate: Voltages → Deflection",
        [
            "Task: Predict trailing-edge deflection from MFC voltages",
            "Best model: Linear regression",
            "Accuracy: R² = 0.9999, MAE = 0.0016 mm",
            "Validation: Parity plot (left), Response curves (right)"
        ],
        image1_path=FIGURES_DIR / "fig02_structural_validation.png",
        image2_path=FIGURES_DIR / "fig03_deflection_response.png",
        img_width=2.1
    )
    
    # Slide 5: Aerodynamic Surrogate Model
    add_two_image_slide(
        prs,
        "Aerodynamic Surrogate: Deflection + Pitch → C_L, C_D",
        [
            "Task: Predict lift and drag from deflection and pitch angle",
            "Baseline: Linear regression (R² = 0.98 for C_L, 0.83 for C_D)",
            "Improved: MLP regressor after augmentation",
            "Results: C_L MAE = 0.0022, C_D MAE = 0.0021 (both R² > 0.99)"
        ],
        image1_path=FIGURES_DIR / "fig04a_aerodynamic_linear.png",
        image2_path=FIGURES_DIR / "fig04b_aerodynamic_mlp.png",
        img_width=2.1
    )
    
    # Slide 6: Data Augmentation and Model Gain
    add_content_slide(
        prs,
        "Data Augmentation and Performance Gain",
        [
            "Generated 7,797 synthetic rows using physics-consistent interpolation",
            "Strict and noisy variants tested for robustness",
            "MLP benefits more from augmentation than linear models",
            "C_D improvement: 20× reduction in error after augmentation"
        ],
        image_path=FIGURES_DIR / "fig06_model_metrics.png",
        image_width=4.5
    )
    
    # Slide 7: Full Chained ML Pipeline
    add_content_slide(
        prs,
        "Full Chained ML Pipeline: Voltages → C_L, C_D",
        [
            "End-to-end: Structural model → Aerodynamic model",
            "Tested on 201 real CFD validation points",
            "Linear chain: C_L MAE = 0.0050, C_D MAE = 0.0093",
            "MLP chain (selected): C_L MAE = 0.0022, C_D MAE = 0.0021"
        ],
        image_path=FIGURES_DIR / "fig05_chain_comparison.png",
        image_width=4.8
    )
    
    # Slide 8: Inverse Voltage Recommendation
    add_content_slide(
        prs,
        "Inverse Voltage Recommendation",
        [
            "Input: Target C_L and pitch angle (with optional C_D constraint)",
            "Method: Grid search (50 V step) + L-BFGS-B continuous refinement",
            "Output: Ranked voltage pairs with predicted aerodynamic coefficients",
            "Use case: Control system for achieving desired lift at any pitch"
        ],
        image_path=FIGURES_DIR / "fig09_inverse_control.png",
        image_width=4.5
    )
    
    # Slide 9: Error Analysis
    add_content_slide(
        prs,
        "Error Analysis: Model Limitations",
        [
            "Chained model accuracy varies with pitch angle",
            "C_L error: Consistent across pitch range",
            "C_D error: Increases slightly at extreme angles",
            "Recommendation: Use model in −4° to +20° pitch range for best accuracy"
        ],
        image_path=FIGURES_DIR / "fig08_error_vs_pitch.png",
        image_width=4.5
    )
    
    # Slide 10: Conclusion and Key Results
    add_content_slide(
        prs,
        "Conclusion: Key Results and Impact",
        [
            "✓ Structural surrogate: Near-perfect accuracy (R² ≈ 0.9999)",
            "✓ Aerodynamic MLP: Significant improvement with augmented training",
            "✓ Full chain: Fast (<60 ms), accurate replacement for simulation",
            "✓ Inverse control: Enables real-time voltage recommendations for target lift"
        ],
        image_path=FIGURES_DIR / "fig06_model_metrics.png",
        image_width=4.5
    )
    
    # Save presentation
    prs.save(str(OUTPUT_FILE))
    print(f"✓ Presentation saved: {OUTPUT_FILE}")

if __name__ == "__main__":
    main()
